"""Assemble a merchant's performance-review deck (.pptx) from model + charts + narrative.

Customer-facing CSM layout, one shared theme for a consistent, polished look & feel:
  Title → Executive Summary (metric cards + bullets) → one slide per KPI
  (header → metric cards → monthly + quarterly charts → CSM analysis) → Notes & Methodology.
No internal QA/data-quality language ever reaches the deck.
"""
from __future__ import annotations

from pathlib import Path
from typing import Optional

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from src.config import settings
from src.presentation import theme
from src.presentation.chart_generator import render_charts
from src.presentation.deck_schema import (
    DeckModel,
    build_deck_model,
    format_amount,
    format_value,
)

_SLIDE_W = Inches(13.333)
_SLIDE_H = Inches(7.5)
_LOGO = settings.PROJECT_ROOT / "assets" / "riskified_logo.png"  # transparent-bg wordmark


def _add_logo(slide, *, left, top, height) -> None:
    """Place the Riskified wordmark (height-scaled, aspect preserved). No-op if missing."""
    if _LOGO.exists():
        slide.shapes.add_picture(str(_LOGO), left, top, height=height)


def _textbox(slide, text, *, left, top, width, height, size=18, bold=False,
             color=theme.PPTX_NAVY, align=PP_ALIGN.LEFT, font=theme.PPTX_BODY_FONT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    para = tf.paragraphs[0]
    para.alignment = align
    run = para.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.name = font
    run.font.color.rgb = RGBColor(*color)
    return box


def _bullets(slide, items, *, left, top, width, height, size=16):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = para.add_run()
        run.text = f"•  {item}"
        run.font.size = Pt(size)
        run.font.name = theme.PPTX_BODY_FONT
        run.font.color.rgb = RGBColor(*theme.PPTX_GREY)
        para.space_after = Pt(6)
    return box


def _band(slide, *, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*color)
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def _card(slide, *, left, top, width, height, label, value,
          sub=None, accent=theme.PPTX_NAVY, sub_color=theme.PPTX_GREY):
    """A rounded metric card: small grey label, large accent value, optional sub line."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*theme.PPTX_LIGHT)
    shape.line.fill.background()
    shape.shadow.inherit = False
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.06)
    tf.margin_bottom = Inches(0.06)

    p0 = tf.paragraphs[0]
    r = p0.add_run(); r.text = label
    r.font.size = Pt(10); r.font.bold = True; r.font.name = theme.PPTX_BODY_FONT
    r.font.color.rgb = RGBColor(*theme.PPTX_GREY)

    p1 = tf.add_paragraph()
    r = p1.add_run(); r.text = value
    r.font.size = Pt(20); r.font.bold = True; r.font.name = theme.PPTX_BODY_FONT
    r.font.color.rgb = RGBColor(*accent)

    if sub:
        p2 = tf.add_paragraph()
        r = p2.add_run(); r.text = sub
        r.font.size = Pt(10); r.font.name = theme.PPTX_BODY_FONT
        r.font.color.rgb = RGBColor(*sub_color)
    return shape


def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _delta_text(change_pct):
    if change_pct is None:
        return None, theme.PPTX_GREY
    arrow = "▲" if change_pct >= 0 else "▼"
    return f"{arrow} {change_pct:+.1f}% since start", theme.PPTX_GREY


def _trend_accent(insight):
    if insight.improving is None:
        return theme.PPTX_NAVY
    return theme.PPTX_TEAL if insight.improving else theme.PPTX_RED


def _title_slide(prs, deck: DeckModel) -> None:
    s = _blank(prs)
    _add_logo(s, left=Inches(0.8), top=Inches(0.65), height=Inches(0.8))

    _textbox(s, deck.merchant_name, left=Inches(0.8), top=Inches(2.7),
             width=Inches(11.7), height=Inches(1.2), size=40, bold=True)
    _textbox(s, "Merchant Performance Review", left=Inches(0.8), top=Inches(3.8),
             width=Inches(11.7), height=Inches(0.8), size=24, color=theme.PPTX_BLUE)
    _band(s, left=Inches(0.85), top=Inches(4.65), width=Inches(2.6), height=Pt(3),
          color=theme.PPTX_TEAL)
    sub = (f"{deck.pre_or_post} authorization   ·   {deck.business_structure}   ·   "
           f"{deck.period_start} – {deck.period_end}")
    _textbox(s, sub, left=Inches(0.8), top=Inches(4.9), width=Inches(11.7),
             height=Inches(0.6), size=16, color=theme.PPTX_GREY)

    _band(s, left=0, top=Inches(6.95), width=_SLIDE_W, height=Inches(0.55), color=theme.PPTX_NAVY)
    _textbox(s, "Prepared by Riskified Customer Success", left=Inches(0.8), top=Inches(7.04),
             width=Inches(11.7), height=Inches(0.4), size=12, color=theme.PPTX_WHITE)


def _kpi_cards_row(slide, kpis, *, top) -> None:
    """A row of one card per KPI showing the latest count value + trend (exec summary)."""
    n = len(kpis)
    gap, left0, total = 0.3, 0.6, 12.13
    cw = (total - gap * (n - 1)) / n
    for i, k in enumerate(kpis):
        left = Inches(left0 + i * (cw + gap))
        sub, sub_color = _delta_text(k.change_pct)
        _card(slide, left=left, top=top, width=Inches(cw), height=Inches(1.15),
              label=k.metric_name, value=format_value(k.unit, k.latest_value),
              sub=sub, accent=_trend_accent(k), sub_color=sub_color)


def _exec_summary_slide(prs, deck, bullets) -> None:
    s = _blank(prs)
    _textbox(s, "Executive Summary", left=Inches(0.6), top=Inches(0.45),
             width=Inches(10), height=Inches(0.9), size=28, bold=True)
    _add_logo(s, left=Inches(12.15), top=Inches(0.35), height=Inches(0.3))
    _kpi_cards_row(s, deck.kpis, top=Inches(1.5))
    _bullets(s, bullets or ["(no summary)"], left=Inches(0.8), top=Inches(3.0),
             width=Inches(11.7), height=Inches(4.0), size=18)


def _kpi_slide(prs, insight, charts, analysis) -> None:
    s = _blank(prs)
    direction = {True: "higher is better", False: "lower is better", None: "transaction volume"}[
        insight.higher_is_better
    ]
    _textbox(s, insight.metric_name, left=Inches(0.6), top=Inches(0.35),
             width=Inches(9), height=Inches(0.7), size=26, bold=True)
    _textbox(s, direction, left=Inches(0.62), top=Inches(1.0), width=Inches(9),
             height=Inches(0.4), size=12, color=theme.PPTX_GREY)
    _add_logo(s, left=Inches(12.15), top=Inches(0.35), height=Inches(0.3))

    # Metric cards: Latest (count), Change, Latest quarter, + Submitted value for Strategic.
    cards = [
        ("Latest", format_value(insight.unit, insight.latest_value),
         insight.latest_period, _trend_accent(insight)),
        ("Change since start",
         "n/a" if insight.change_pct is None else f"{insight.change_pct:+.1f}%",
         "vs first month", _trend_accent(insight)),
        ("Latest quarter",
         format_value(insight.unit, insight.quarterly_values[-1]) if insight.quarterly_values else "n/a",
         insight.quarterly_periods[-1] if insight.quarterly_periods else "", theme.PPTX_NAVY),
    ]
    if insight.amount is not None:
        a = insight.amount
        # Submission Volume's amount view is a dollar figure; rate KPIs' is a percentage.
        amount_label = "Submitted value" if insight.metric_id == "submission_volume" \
            else "Amount-weighted latest"
        cards.append((
            amount_label,
            format_amount(insight.metric_id, a.unit, a.latest_value),
            a.latest_period or "", theme.PPTX_BLUE,
        ))

    n = len(cards)
    gap, left0, total = 0.3, 0.6, 12.13
    cw = (total - gap * (n - 1)) / n
    for i, (label, value, sub, accent) in enumerate(cards):
        _card(s, left=Inches(left0 + i * (cw + gap)), top=Inches(1.45),
              width=Inches(cw), height=Inches(1.0), label=label, value=value,
              sub=sub, accent=accent)

    s.shapes.add_picture(str(charts["monthly"]), Inches(0.5), Inches(2.75), width=Inches(5.9))
    s.shapes.add_picture(str(charts["quarterly"]), Inches(6.95), Inches(2.75), width=Inches(5.9))
    _textbox(s, analysis or "", left=Inches(0.6), top=Inches(5.95),
             width=Inches(12.1), height=Inches(1.4), size=13, color=theme.PPTX_GREY)


def _methodology_bullets(deck: DeckModel) -> list[str]:
    """Short, customer-safe appendix. Enterprise decks never mention the amount-weighted view."""
    strategic = any(k.amount is not None for k in deck.kpis)
    if strategic:
        bullets = [
            "Count-based view: KPIs by transaction volume (number of orders).",
            "Amount-weighted view: KPIs weighted by submitted order value "
            "(Submission Volume is shown in dollars).",
        ]
    else:
        bullets = [
            "This deck uses the count-based (transaction-volume) view.",
        ]
    bullets.append("Charts annotate notable evidence events from the period.")
    return bullets


def _notes_slide(prs, deck: DeckModel) -> None:
    s = _blank(prs)
    _textbox(s, "Notes & Methodology", left=Inches(0.6), top=Inches(0.5),
             width=Inches(10), height=Inches(0.9), size=28, bold=True)
    _add_logo(s, left=Inches(12.15), top=Inches(0.35), height=Inches(0.3))
    _bullets(s, _methodology_bullets(deck), left=Inches(0.8), top=Inches(1.7),
             width=Inches(11.7), height=Inches(5), size=16)


def build_deck(deck: DeckModel, charts, narrative, out_path: Path | str) -> Path:
    """Assemble the .pptx and write it to the explicit ``out_path``."""
    prs = Presentation()
    prs.slide_width, prs.slide_height = _SLIDE_W, _SLIDE_H

    _title_slide(prs, deck)
    _exec_summary_slide(prs, deck, narrative.executive_summary)
    for insight in deck.kpis:
        _kpi_slide(prs, insight, charts[insight.metric_id],
                   narrative.kpi_analysis.get(insight.metric_id, ""))
    _notes_slide(prs, deck)

    out_path = Path(out_path)
    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    return out_path


def generate_deck_for_merchant(
    repo,
    merchant_id: str,
    *,
    model: Optional[object] = None,
    version: Optional[str] = None,
    base: Path | str = settings.OUTPUT_DIR,
) -> Path:
    """Full per-merchant deck pipeline: model → narrative → **eval gate** → versioned .pptx.

    Narrative evaluation is a real gate: if it fails, no charts/pptx are written and
    ``LATEST`` is NOT updated — a ``NarrativeEvalError`` is raised with the reason. On
    success, writes the deck + a ``<deck>.eval.json`` sidecar and refreshes both LATEST
    pointers at ``base/decks/<merchant_id>/<version>.pptx`` and
    ``base/charts/<merchant_id>/<version>/``.
    """
    from src.llm.evaluation import NarrativeEvalError, evaluate_narrative
    from src.llm.narrative_generator import generate_narrative
    from src.presentation.versioning import merchant_paths, new_version, write_latest

    version = version or new_version()
    deck = build_deck_model(repo, merchant_id)
    narrative = generate_narrative(deck, model=model)

    # GATE: a failing narrative blocks the deck before anything is written.
    report = evaluate_narrative(deck, narrative)
    if not report.ok:
        raise NarrativeEvalError(f"{merchant_id}: narrative eval failed — {report.reason()}")

    deck_path, charts_dir = merchant_paths(merchant_id, version, base)
    charts = render_charts(deck, charts_dir)
    path = build_deck(deck, charts, narrative, deck_path)
    (path.parent / f"{path.stem}.eval.json").write_text(report.model_dump_json(indent=2))

    write_latest(Path(base) / "decks" / merchant_id, version)
    write_latest(Path(base) / "charts" / merchant_id, version)
    return path
