"""Deck assembly — build a .pptx end-to-end with a fake LLM, then re-open it."""
from pptx import Presentation

from src.db.connection import get_connection
from src.ingestion.loaders import read_kpis, read_profiles, read_evidence
from src.llm.narrative_generator import NarrativeBundle
from src.metrics.registry import METRIC_REGISTRY
from src.pipeline import run_pipeline
from src.presentation.deck_generator import generate_deck_for_merchant
from src.repositories.metrics_repository import MetricsRepository


class _FakeStructured:
    def __init__(self, result):
        self.result = result

    def invoke(self, messages):
        return self.result


class FakeChatModel:
    def __init__(self, result):
        self.result = result

    def with_structured_output(self, schema):
        return _FakeStructured(self.result)


def _repo():
    con = get_connection(":memory:")
    run_pipeline(con, read_kpis(), read_profiles(), read_evidence())
    return MetricsRepository(con)


def _fake_model():
    return FakeChatModel(NarrativeBundle(
        executive_summary=["Approval rate is healthy.", "Volume steady through the year."],
        kpi_analysis={m.id: f"Analysis for {m.name}." for m in METRIC_REGISTRY},
    ))


def test_generate_deck_creates_pptx_with_expected_slides(tmp_path):
    path = generate_deck_for_merchant(
        _repo(), "acme", model=_fake_model(), version="20260101T000000Z", base=tmp_path,
    )
    assert path == tmp_path / "decks" / "acme" / "acme_20260101T000000Z.pptx"
    assert path.exists() and path.suffix == ".pptx"

    prs = Presentation(str(path))
    # Title + Executive Summary + 4 KPI slides + Data-Quality Notes = 7
    assert len(prs.slides) == 7

    def slide_text(slide):
        return "\n".join(s.text for s in slide.shapes if s.has_text_frame)

    assert "ACME" in slide_text(prs.slides[0])
    assert "Executive Summary" in slide_text(prs.slides[1])
    all_text = "\n".join(slide_text(s) for s in prs.slides)
    assert "Analysis for Approval Rate." in all_text
    assert "Notes & Methodology" in all_text


def test_layout_is_company_then_version_with_latest_pointer(tmp_path):
    path = generate_deck_for_merchant(
        _repo(), "vandelay-industries", model=_fake_model(), version="v1", base=tmp_path,
    )
    assert path.name == "vandelay-industries_v1.pptx"
    assert path.parent == tmp_path / "decks" / "vandelay-industries"
    # charts live under charts/<merchant>/<version>/
    assert (tmp_path / "charts" / "vandelay-industries" / "v1").is_dir()
    # per-merchant LATEST pointer
    assert (tmp_path / "decks" / "vandelay-industries" / "LATEST").read_text().strip() == "v1"
    # advisory evaluation sidecar written next to the deck
    import json
    sidecar = tmp_path / "decks" / "vandelay-industries" / "vandelay-industries_v1.eval.json"
    assert sidecar.exists()
    report = json.loads(sidecar.read_text())
    assert report["merchant_id"] == "vandelay-industries"
    assert report["ok"] is True  # fake narrative covers all KPIs, no invented numbers


def test_eval_failure_gates_the_deck(tmp_path):
    import pytest
    from src.llm.evaluation import NarrativeEvalError

    # Narrative missing a KPI analysis -> structural failure -> deck must NOT be produced.
    broken = FakeChatModel(NarrativeBundle(
        executive_summary=["Summary."],
        kpi_analysis={m.id: "x" for m in METRIC_REGISTRY if m.id != "effective_fraud_rate"},
    ))
    with pytest.raises(NarrativeEvalError):
        generate_deck_for_merchant(_repo(), "acme", model=broken, version="v1", base=tmp_path)

    # No pptx and no LATEST were written for the failed merchant.
    assert not (tmp_path / "decks" / "acme" / "acme_v1.pptx").exists()
    assert not (tmp_path / "decks" / "acme" / "LATEST").exists()


# --- customer-facing deck v2 ---

def _all_text(path):
    prs = Presentation(str(path))
    return "\n".join(
        s.text for slide in prs.slides for s in slide.shapes if s.has_text_frame
    )


def _slides_text(path):
    """Per-slide text. KPI slides start at index 2 (after Title + Exec Summary)."""
    prs = Presentation(str(path))
    return ["\n".join(s.text for s in slide.shapes if s.has_text_frame) for slide in prs.slides]


def _kpi_slide_text(path):
    """metric_id -> that KPI slide's text (slides 2.. follow METRIC_REGISTRY order)."""
    slides = _slides_text(path)
    return {m.id: slides[2 + i] for i, m in enumerate(METRIC_REGISTRY)}


def test_deck_has_no_internal_qa_language(tmp_path):
    from src.llm.evaluation import _BANNED_CAUSAL, _BANNED_INTERNAL
    path = generate_deck_for_merchant(_repo(), "acme", model=_fake_model(),
                                      version="v1", base=tmp_path)
    text = _all_text(path).lower()
    for term in _BANNED_INTERNAL + _BANNED_CAUSAL:
        assert term not in text, f"internal/QA term leaked into deck: {term!r}"
    assert "data points" not in text


def test_strategic_deck_shows_amount_enterprise_count_only(tmp_path):
    acme = _all_text(generate_deck_for_merchant(
        _repo(), "acme", model=_fake_model(), version="v1", base=tmp_path))
    vandelay = _all_text(generate_deck_for_merchant(
        _repo(), "vandelay-industries", model=_fake_model(), version="v1", base=tmp_path))
    # Strategic surfaces the amount (submitted value) card; Enterprise does not.
    assert "Submitted value" in acme
    assert "Submitted value" not in vandelay


def test_notes_slide_is_customer_safe_methodology(tmp_path):
    text = _all_text(generate_deck_for_merchant(
        _repo(), "acme", model=_fake_model(), version="v1", base=tmp_path))
    assert "Notes & Methodology" in text
    assert "transaction volume" in text.lower()
    assert "weighted by submitted order value" in text


def test_enterprise_methodology_omits_amount_weighted(tmp_path):
    text = _all_text(generate_deck_for_merchant(
        _repo(), "vandelay-industries", model=_fake_model(), version="v1", base=tmp_path))
    assert "count-based" in text.lower()
    assert "amount-weighted" not in text.lower()  # never mention it for Enterprise


def test_each_kpi_slide_embeds_its_own_metric_charts(tmp_path):
    # Guards against wiring the wrong chart image onto a slide (req: correct chart per KPI).
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    path = generate_deck_for_merchant(_repo(), "acme", model=_fake_model(),
                                      version="v1", base=tmp_path)
    charts_dir = tmp_path / "charts" / "acme" / "v1"
    prs = Presentation(str(path))
    for i, m in enumerate(METRIC_REGISTRY):
        slide = prs.slides[2 + i]  # KPI slides follow Title + Exec Summary
        blobs = {sh.image.blob for sh in slide.shapes
                 if sh.shape_type == MSO_SHAPE_TYPE.PICTURE}
        assert (charts_dir / f"{m.id}_monthly.png").read_bytes() in blobs, m.id
        assert (charts_dir / f"{m.id}_quarterly.png").read_bytes() in blobs, m.id
        # and NOT a different metric's chart
        other = METRIC_REGISTRY[(i + 1) % len(METRIC_REGISTRY)].id
        assert (charts_dir / f"{other}_monthly.png").read_bytes() not in blobs, m.id


def test_amount_card_labels_are_semantically_correct(tmp_path):
    by_metric = _kpi_slide_text(generate_deck_for_merchant(
        _repo(), "acme", model=_fake_model(), version="v1", base=tmp_path))
    # Submission Volume: a "Submitted value" card with a dollar figure.
    assert "Submitted value" in by_metric["submission_volume"]
    assert "$" in by_metric["submission_volume"]
    # Rate KPIs: "Amount-weighted latest" (percentage), never "Submitted value".
    for mid in ("approval_rate", "accepted_chargeback_rate", "effective_fraud_rate"):
        assert "Amount-weighted latest" in by_metric[mid], mid
        assert "Submitted value" not in by_metric[mid], mid
        assert "$" not in by_metric[mid], mid


def test_enterprise_deck_has_no_amount_cards(tmp_path):
    text = _all_text(generate_deck_for_merchant(
        _repo(), "vandelay-industries", model=_fake_model(), version="v1", base=tmp_path))
    assert "Amount-weighted latest" not in text
    assert "Submitted value" not in text


def test_deck_embeds_logo(tmp_path):
    # The real logo (PNG asset) is embedded as a picture on the title slide.
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    path = generate_deck_for_merchant(_repo(), "acme", model=_fake_model(),
                                      version="v1", base=tmp_path)
    title = Presentation(str(path)).slides[0]
    assert any(sh.shape_type == MSO_SHAPE_TYPE.PICTURE for sh in title.shapes), \
        "title slide should embed the Riskified logo image"
